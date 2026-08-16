#!/usr/bin/env python
"""Pitch deck builder (python-pptx) -- title + 9 content slides, per docs/PREZENTACIJA_BRIEF.md.

Visual redesign round (user request, day 5):
  - Times New Roman everywhere, no text below 30 pt (small captions reformulated
    or cut -- the speakers narrate them), footers removed
  - per-slide two-stop gradient backgrounds, multi-color accents, gradient bars
  - title slide without mentors
  - slide 3 is a full-slide implementation-architecture diagram (no headline):
    video lane and action lane converge into the DiT, output decoded below
  - slide 4 reduced to the two design decisions, one card each
  - slide 6 ladder rebuilt label-above-bar (nothing side-by-side can overlap),
    one color per rung

Content still follows the brief: 99.6% (not the 100% subset), the ladder,
the click-to-play widget (demo1, scene 204), inverse dynamics 15/20, the
corrected rollout story (control survives, localization drifts), four
contributions. Numbers all from the 256-scene eval unless said otherwise.

Prerequisites:
  - widget tiles/posters in /home/mls10/presentation/widget (extract_demo1_widget.py)
  - goal-search panel: logs/goal_search/goal_search_idx108.png
  - cp logs/demo/index.html /home/mls10/presentation/arm_control_panel.html

Round 2 (user request): the architecture diagram's click-to-reveal highlight
(raw OOXML <p:timing>, unverifiable in this container) was replaced by a
second static slide -- architecture_slide(prs, highlight=True/False), called
twice. The controllability ladder now wraps each bar pair in its own
translucent, color-matched bubble with a pill tab naming the metric (video
quality vs. action accuracy). The limitations slide's cards were relabeled
("control holds" / "video quality decay") and its background switched to
solid (new_slide_solid) to drop the diagonal gradient seam that cut across
the three edge-to-edge cards.

Gotchas: python-pptx 1.0.2 dash style lives in pptx.enum.dml.MSO_LINE_DASH_STYLE;
hyperlink text must live in the linked shape's own text_frame; connector
arrowheads need raw oxml (a:tailEnd).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette --
INK    = RGBColor(0xF6, 0xF7, 0xFB)
DIM    = RGBColor(0xCB, 0xD3, 0xE4)
PANEL  = RGBColor(0x10, 0x14, 0x26)
TEAL   = RGBColor(0x3E, 0xE6, 0xD0)
PINK   = RGBColor(0xFF, 0x7A, 0xB8)
AMBER  = RGBColor(0xFF, 0xD1, 0x66)
CYAN   = RGBColor(0x6F, 0xD3, 0xFF)
VIOLET = RGBColor(0xB7, 0x9B, 0xFF)
CORAL  = RGBColor(0xFF, 0x8F, 0x6B)
SLATE  = RGBColor(0x9A, 0xA6, 0xBE)

FONT = "Times New Roman"

WIDGET = "/home/mls10/presentation/widget"
GOAL_PANEL = "/home/mls10/logs/goal_search/goal_search_idx108.png"
OUT_PATH = "/home/mls10/presentation/BAIR_LoRA_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ------------------------------------------------------------- utilities --
def set_letter_spacing(run, pt_hundredths):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(pt_hundredths))


def grad(fill, c1, c2, angle=120.0):
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    fill.gradient_angle = angle


def new_slide(prs, c1, c2, angle=125.0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    grad(slide.background.fill, c1, c2, angle)
    return slide


def new_slide_solid(prs, color):
    """Flat background, no diagonal gradient -- for slides where the usual
    120ish-degree two-stop gradient reads as an unwanted diagonal seam."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    return slide


def add_text(slide, left, top, width, height, text, size, color,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.08, letter_spacing=None,
             upper=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line.upper() if upper else line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT
        run.font.color.rgb = color
        if letter_spacing is not None:
            set_letter_spacing(run, letter_spacing)
    return tb


def eyebrow(slide, left, top, text, color, align=PP_ALIGN.LEFT, width=Inches(12)):
    add_text(slide, left, top, width, Inches(0.55), text, 30, color,
             bold=True, align=align, letter_spacing=120, upper=True)


def rect(slide, left, top, width, height, fill_color=None, line_color=None,
         radius=False, line_w=1.5, grad_pair=None, grad_angle=90.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if grad_pair is not None:
        grad(shp.fill, grad_pair[0], grad_pair[1], grad_angle)
    elif fill_color is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.10
        except Exception:
            pass
    return shp


def accent_bar(slide, left, top, width, c1, c2, height=Pt(4)):
    bar = rect(slide, left, top, width, height, grad_pair=(c1, c2), grad_angle=0.0)
    return bar


def box(slide, left, top, width, height, lines, border, size=30,
        fill=PANEL, text_color=INK):
    rect(slide, left, top, width, height, fill_color=fill, line_color=border,
         radius=True, line_w=2.25)
    add_text(slide, left + Inches(0.08), top, width - Inches(0.16), height,
             lines, size, text_color, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def block_arrow_right(slide, left, top, width, height, color, grad_pair=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    if grad_pair is not None:
        grad(shp.fill, grad_pair[0], grad_pair[1], 0.0)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def arrow_conn(slide, x1, y1, x2, y2, color, w=3.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "lg", "len": "lg"})
    ln.append(tail)
    return conn


def chip(slide, left, top, width, height, text, border, size=30):
    rect(slide, left, top, width, height, fill_color=PANEL, line_color=border,
         radius=True, line_w=2.0)
    add_text(slide, left + Inches(0.1), top, width - Inches(0.2), height,
             text, size, INK, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def set_fill_alpha(shape, pct):
    """Set solid-fill opacity (0-100). Call after fill.solid()+fore_color.rgb."""
    srgb = shape.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    srgb.append(alpha)


def stage_band(slide, left, top, width, height, color, alpha_pct=14):
    """Translucent full-height rounded band marking one vertical pipeline stage."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    set_fill_alpha(shp, alpha_pct)
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def ring(slide, box_left, box_top, box_w, box_h, pad, color, line_w=3.25):
    """Dashed highlight oval, generously padded so it fully contains the box."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, box_left - pad, box_top - pad,
                                 box_w + 2 * pad, box_h + 2 * pad)
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(line_w)
    shp.line.dash_style = MSO_LINE.DASH
    shp.shadow.inherit = False
    return shp


def architecture_slide(prs, highlight):
    """Pipeline-stage architecture diagram. Built twice (highlight=False then
    True) instead of using a PowerPoint click-animation for the highlight
    rings -- this container has no PowerPoint/LibreOffice to render OOXML
    timing XML, so a hand-built <p:timing> block could not be verified and
    was unreliable in practice. Two static slides always render correctly."""
    s = new_slide(prs, RGBColor(0x12, 0x3A, 0x5C), RGBColor(0x3A, 0x1D, 0x5E), 115)

    STAGE_INPUT, STAGE_ENCODE, STAGE_READY, STAGE_FUSE = CYAN, AMBER, PINK, VIOLET
    NEW_COLOR = TEAL   # marks the two pieces we integrated on top of the backbone

    lane_h = Inches(1.3)
    y1 = Inches(1.2)     # video lane (top)
    y2 = Inches(5.35)    # action lane (bottom)

    # -- stage bands first (z-order: behind everything else) ----------------
    band_top, band_bottom = Inches(0.95), Inches(6.9)
    band_h = band_bottom - band_top
    stage_band(s, Inches(0.40), band_top, Inches(2.90), band_h, STAGE_INPUT)
    stage_band(s, Inches(4.10), band_top, Inches(2.55), band_h, STAGE_ENCODE)
    stage_band(s, Inches(7.45), band_top, Inches(2.45), band_h, STAGE_READY)
    stage_band(s, Inches(10.05), Inches(0.35), Inches(2.85), Inches(4.60), STAGE_FUSE)

    # video lane
    box(s, Inches(0.6), y1, Inches(2.5), lane_h, "BAIR video\n64 × 64", STAGE_INPUT)
    block_arrow_right(s, Inches(3.3), y1 + Inches(0.4), Inches(0.85), Inches(0.5), STAGE_ENCODE)
    box(s, Inches(4.35), y1, Inches(2.1), lane_h, "VAE\nencode", STAGE_ENCODE)
    block_arrow_right(s, Inches(6.65), y1 + Inches(0.4), Inches(0.85), Inches(0.5), STAGE_READY)
    box(s, Inches(7.7), y1, Inches(2.1), lane_h, "4 latent\nframes", STAGE_READY)

    # action lane
    box(s, Inches(0.6), y2, Inches(2.5), lane_h, "actions\n4 × 4D", STAGE_INPUT)
    block_arrow_right(s, Inches(3.3), y2 + Inches(0.4), Inches(0.85), Inches(0.5), STAGE_ENCODE)
    box(s, Inches(4.35), y2, Inches(2.1), lane_h, "MLP\n16 → 1536", STAGE_ENCODE)
    block_arrow_right(s, Inches(6.65), y2 + Inches(0.4), Inches(0.85), Inches(0.5), STAGE_READY)
    box(s, Inches(7.7), y2, Inches(2.1), lane_h, "+ timestep\nAdaLN", STAGE_READY)

    # the DiT core (center right, violet gradient -- the fuse stage)
    rect(s, Inches(10.15), Inches(2.75), Inches(2.55), Inches(2.0),
         grad_pair=(RGBColor(0x6C, 0x3F, 0xB8), RGBColor(0x2E, 0x7E, 0x8C)),
         line_color=STAGE_FUSE, radius=True, line_w=2.5, grad_angle=35.0)
    add_text(s, Inches(10.15), Inches(2.75), Inches(2.55), Inches(2.0),
             "30 DiT\nblocks\n+ LoRA r16", 30, INK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)

    # converging arrows into the DiT, and out to the decoded video
    arrow_conn(s, Inches(9.8), y1 + Inches(0.65), Inches(10.35), Inches(3.05), STAGE_FUSE)
    arrow_conn(s, Inches(9.8), y2 + Inches(0.65), Inches(10.35), Inches(4.45), STAGE_FUSE)
    box(s, Inches(10.6), Inches(0.55), Inches(2.1), Inches(1.1), "next 4\nlatents", STAGE_FUSE)
    arrow_conn(s, Inches(11.65), Inches(2.75), Inches(11.65), Inches(1.7), STAGE_FUSE)

    add_text(s, Inches(0.6), Inches(3.35), Inches(8.5), Inches(0.7),
             "1.4 % trainable  ·  15 lines changed upstream", 30, AMBER, bold=True)

    # -- stage legend, bottom strip -------------------------------------------
    legend = [("input", STAGE_INPUT, Inches(0.6), Inches(1.3)),
              ("encode", STAGE_ENCODE, Inches(2.55), Inches(1.5)),
              ("ready", STAGE_READY, Inches(4.75), Inches(1.3)),
              ("fuse + decode", STAGE_FUSE, Inches(6.75), Inches(3.0))]
    for name, col, x, w in legend:
        rect(s, x, Inches(7.05), Inches(0.34), Inches(0.34), fill_color=col, radius=True)
        add_text(s, x + Inches(0.46), Inches(6.98), w, Inches(0.48),
                 name, 30, col, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    if highlight:
        # second pass of the same slide: what we actually added on top of
        # the pretrained backbone, circled and labeled -- static, so it
        # always renders correctly instead of depending on a click animation
        chip(s, Inches(0.6), Inches(0.42), Inches(3.7), Inches(0.55),
             "what we added ↓", NEW_COLOR, 30)
        ring(s, Inches(4.35), y2, Inches(2.1), lane_h, Inches(0.4), NEW_COLOR)
        ring(s, Inches(10.15), Inches(2.75), Inches(2.55), Inches(2.0), Inches(0.5), NEW_COLOR)

    return s


# ------------------------------------------------------------------ build --
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================ 1 · Title ====
s = new_slide(prs, RGBColor(0x2B, 0x1A, 0x5E), RGBColor(0x0B, 0x3B, 0x43), 130)

# PSIML logo, top-left -- background is now transparent (prep_psiml_logo.py),
# so it drops onto the gradient cleanly instead of showing a mismatched navy
# box; nothing else shares this row, so no collision risk
logo_w, logo_h = Inches(1.9), Inches(1.9 * 251 / 603)
s.shapes.add_picture("/home/mls10/presentation/psiml_logo.png",
                     Inches(0.55), Inches(0.42), logo_w, logo_h)

add_text(s, Inches(1), Inches(1.75), Inches(11.33), Inches(0.6),
         "BAIR · WAN2.1-1.3B · LORA", 30, CYAN, bold=True,
         align=PP_ALIGN.CENTER, letter_spacing=140)
add_text(s, Inches(1), Inches(2.45), Inches(11.33), Inches(2.0),
         "Action-Conditioned\nWorld Models", 54, INK, bold=True,
         align=PP_ALIGN.CENTER, line_spacing=1.02)
accent_bar(s, Inches(5.27), Inches(4.55), Inches(2.8), PINK, TEAL, Pt(5))
add_text(s, Inches(1.4), Inches(4.95), Inches(10.53), Inches(0.75),
         "Are world models all we need?", 32, DIM, italic=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.35), Inches(11.33), Inches(0.6),
         "David Marković  &  Mihajlo Stevanović", 30, AMBER, bold=True,
         align=PP_ALIGN.CENTER)
# mentor credit, deliberately smaller -- a byline under the authors, not a
# second headline competing with them
add_text(s, Inches(1), Inches(7.0), Inches(11.33), Inches(0.4),
         "Mentors:  Nedko Savov  ·  Danilo Đorđević  ·  Nataša Jovanović",
         18, DIM, align=PP_ALIGN.CENTER)

# ======================================================= 2 · Motivation ====
s = new_slide(prs, RGBColor(0x1C, 0x2A, 0x6B), RGBColor(0x0E, 0x4A, 0x4A), 120)
eyebrow(s, Inches(0.8), Inches(0.5), "Motivation", TEAL)
add_text(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.6),
         "A robot that can imagine what happens next\ncan plan before it acts",
         38, INK, bold=True, line_spacing=1.1)
accent_bar(s, Inches(0.8), Inches(2.95), Inches(3.2), TEAL, CYAN)

add_text(s, Inches(0.8), Inches(3.45), Inches(5.9), Inches(1.3),
         "Without a world model,\nyou learn the hard way",
         32, INK, bold=True, line_spacing=1.15)
for i, (g, c) in enumerate([("action control", PINK),
                            ("precision", AMBER),
                            ("safety", TEAL)]):
    chip(s, Inches(0.8), Inches(4.85) + i * Inches(0.87), Inches(4.9),
         Inches(0.75), g, c, 30)

# Click-to-play fail clip (Unitree H1 flailing on its crane, 12.5 s, muted;
# cut at 1:54 from youtube OqsqzPCLvE4 per user request -- talk use only, the
# compilation is not ours to redistribute). The DRC Atlas fall photo
# (robot_fall.jpg, CC BY 2.0) stays on disk as a swap-in alternative.
img_w, img_h = Inches(6.2), Inches(3.49)   # clip is 640x360, keep 16:9
ix, iy = Inches(6.15), Inches(3.35)
rect(s, ix - Inches(0.05), iy - Inches(0.05), img_w + Inches(0.10),
     img_h + Inches(0.10), fill_color=PANEL, line_color=CORAL, radius=False,
     line_w=2.5)
s.shapes.add_movie("/home/mls10/presentation/robot_fail_clip.mp4", ix, iy,
                   img_w, img_h,
                   poster_frame_image="/home/mls10/presentation/robot_fail_poster.png",
                   mime_type="video/mp4")

# ==================================== 3a · High-level pipeline (real data) ==
# Frames/clip are our own quickstart generation (cheeseburger scene, camera
# command forward x19). Narrate that this is the pretrained minWM demo; the
# BAIR fine-tune runs the same loop at 64x64.
s = new_slide(prs, RGBColor(0x0E, 0x46, 0x3E), RGBColor(0x24, 0x1B, 0x4F), 120)
add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.75),
         "One frame and a command in, a video out", 36, INK, bold=True)
accent_bar(s, Inches(0.7), Inches(1.2), Inches(3.2), TEAL, VIOLET)

# input frame + action, each its own structure with its own arrow into the
# model (the merge belongs one stage later -- model + latent tokens)
rect(s, Inches(0.55), Inches(1.85), Inches(2.7), Inches(1.6),
     fill_color=PANEL, line_color=CYAN, radius=False, line_w=2.5)
s.shapes.add_picture("/home/mls10/presentation/burger_in.png",
                     Inches(0.6), Inches(1.9), Inches(2.6), Inches(1.5))
add_text(s, Inches(0.5), Inches(3.55), Inches(2.8), Inches(0.42),
         "current frame", 30, CYAN, bold=True, align=PP_ALIGN.CENTER)
chip(s, Inches(0.6), Inches(4.6), Inches(2.6), Inches(0.75), "forward ↥", PINK, 30)
add_text(s, Inches(0.5), Inches(5.5), Inches(2.8), Inches(0.42),
         "action", 30, PINK, bold=True, align=PP_ALIGN.CENTER)

block_arrow_right(s, Inches(3.35), Inches(2.4), Inches(0.6), Inches(0.5), CYAN)
block_arrow_right(s, Inches(3.35), Inches(4.72), Inches(0.6), Inches(0.5), PINK)

# model + latent tokens are now ONE block: the model produces the tokens, so
# they share a single outer panel instead of being connected by an arrow
rect(s, Inches(4.15), Inches(2.15), Inches(4.9), Inches(3.2),
     grad_pair=(RGBColor(0x6C, 0x3F, 0xB8), RGBColor(0x2E, 0x7E, 0x8C)),
     line_color=AMBER, radius=True, line_w=2.5, grad_angle=35.0)
add_text(s, Inches(4.35), Inches(2.55), Inches(2.3), Inches(1.0),
         "video\nmodel", 34, INK, bold=True, align=PP_ALIGN.CENTER,
         line_spacing=1.05)
add_text(s, Inches(4.35), Inches(4.15), Inches(2.3), Inches(0.6),
         "block by block", 30, INK, align=PP_ALIGN.CENTER)

# latent-token grid, inset on its own dark backing for contrast against the
# model's gradient panel (no border of its own -- reads as part of one block)
rect(s, Inches(6.95), Inches(2.5), Inches(1.98), Inches(1.9),
     fill_color=PANEL, radius=True)
tok_cols = [TEAL, CYAN, VIOLET, PINK]
for r_ in range(4):
    for c_ in range(4):
        rect(s, Inches(6.95) + Inches(0.19) + c_ * Inches(0.40),
             Inches(2.5) + Inches(0.19) + r_ * Inches(0.40),
             Inches(0.32), Inches(0.32), fill_color=tok_cols[(r_ + c_) % 4],
             radius=True)
add_text(s, Inches(6.65), Inches(4.55), Inches(2.2), Inches(0.55),
         "latent tokens", 30, VIOLET, bold=True, align=PP_ALIGN.CENTER)

block_arrow_right(s, Inches(9.15), Inches(3.35), Inches(0.4), Inches(0.5), AMBER)
# sits in the open band above the boxes -- neither neighboring box starts
# higher than y=2.15, so this clears both with room to spare
add_text(s, Inches(8.65), Inches(1.6), Inches(1.6), Inches(0.45),
         "decode", 30, AMBER, bold=True, align=PP_ALIGN.CENTER)

# output: the playable generated clip -- enlarged again (was 2.9x1.674in)
mv_w, mv_h = Inches(3.43), Inches(1.98)   # 832x480, keep 16:9
mx, my = Inches(9.70), Inches(2.45)
rect(s, mx - Inches(0.05), my - Inches(0.05), mv_w + Inches(0.10),
     mv_h + Inches(0.10), fill_color=PANEL, line_color=TEAL, radius=False,
     line_w=2.5)
s.shapes.add_movie("/home/mls10/presentation/burger_demo.mp4", mx, my, mv_w, mv_h,
                   poster_frame_image="/home/mls10/presentation/burger_mid1.png",
                   mime_type="video/mp4")
add_text(s, Inches(9.45), Inches(4.6), Inches(3.7), Inches(0.55),
         "generated video ▶", 30, TEAL, bold=True, align=PP_ALIGN.CENTER)

# ============================================ 3 · Architecture diagram =====
# Colored by PIPELINE STAGE (a vertical band per column), not by modality/lane
# -- both the video-lane and action-lane box in a given column share one
# color, so the whole vertical slice reads as one step of the pipeline.
# Shown as two slides: plain, then the same diagram again with the two
# pieces we added circled (see architecture_slide() -- was a click animation,
# swapped for a second static slide, see its docstring for why).
architecture_slide(prs, highlight=False)
architecture_slide(prs, highlight=True)

# ====================================== 4 · The two design decisions =======
s = new_slide(prs, RGBColor(0x3D, 0x17, 0x48), RGBColor(0x14, 0x35, 0x5C), 120)
eyebrow(s, Inches(0.8), Inches(0.5), "Method", VIOLET)
add_text(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.9),
         "Two architecture pieces", 40, INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1), Inches(3.2), PINK, VIOLET)

cw, ch2 = Inches(5.75), Inches(3.7)
cy = Inches(2.6)
# piece 1: adapts the PRETRAINED backbone's own weights
rect(s, Inches(0.8), cy, cw, ch2, fill_color=PANEL, line_color=PINK,
     radius=True, line_w=2.5)
add_text(s, Inches(1.15), cy + Inches(0.35), cw - Inches(0.7), Inches(0.7),
         "LoRA: QKV + FFN", 36, PINK, bold=True, line_spacing=1.05)
add_text(s, Inches(1.15), cy + Inches(1.25), cw - Inches(0.7), Inches(2.2),
         "rank 16, every block\n~1.4% trainable\nadapts existing weights",
         30, DIM, line_spacing=1.25)

# piece 2: the genuinely NEW pathway -- this is what lets the model see the
# action at all, LoRA alone changes nothing about that
rect(s, Inches(6.85), cy, cw, ch2, fill_color=PANEL, line_color=TEAL,
     radius=True, line_w=2.5)
add_text(s, Inches(7.2), cy + Inches(0.35), cw - Inches(0.7), Inches(0.7),
         "Action Encoder", 36, TEAL, bold=True, line_spacing=1.05)
add_text(s, Inches(7.2), cy + Inches(1.25), cw - Inches(0.7), Inches(2.2),
         "16 → 256 → 256\n→ 1536, zero-init\nflatten, never average\n→ AdaLN, every frame",
         30, DIM, line_spacing=1.1)

# ================================================= 5 · Result: control ====
s = new_slide(prs, RGBColor(0x0F, 0x4D, 0x43), RGBColor(0x1A, 0x1F, 0x5C), 125)
eyebrow(s, Inches(0), Inches(0.6), "Result", TEAL, align=PP_ALIGN.CENTER,
        width=Inches(13.333))
add_text(s, Inches(0), Inches(1.05), Inches(13.333), Inches(2.5),
         "99.6%", 150, INK, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.2), Inches(3.75), Inches(10.93), Inches(0.7),
         "of held-out scenes follow the command  (255 / 256)",
         32, DIM, align=PP_ALIGN.CENTER)

cw, ch, cgap = Inches(2.55), Inches(0.8), Inches(0.25)
row_w = cw * 4 + cgap * 3
cx = (SLIDE_W - row_w) / 2
for glyph, c in [("↑ up ✓", CYAN), ("↓ down ✓", VIOLET),
                 ("← left ✓", PINK), ("→ right ✓", TEAL)]:
    chip(s, cx, Inches(4.75), cw, ch, glyph, c, 30)
    cx += cw + cgap

add_text(s, Inches(1.2), Inches(5.85), Inches(10.93), Inches(0.6),
         "absolute position accuracy: 84.8%", 30, AMBER, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(6.55), Inches(10.93), Inches(0.6),
         "PSNR 18.6 dB   ·   SSIM 0.79   ·   FID 11.1", 30, DIM,
         align=PP_ALIGN.CENTER)

# ============================================ 6 · Controllability ladder ====
s = new_slide(prs, RGBColor(0x35, 0x19, 0x4F), RGBColor(0x0C, 0x4A, 0x5E), 118)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9),
         "What is the tuning worth?", 40, INK, bold=True)
accent_bar(s, Inches(0.8), Inches(1.45), Inches(3.2), AMBER, PINK)

chart_left = Inches(0.8)
chart_width = Inches(10.6)
scale_max = 24.0
ceiling_db = 22.74
# top pair = video quality (PSNR), bottom pair = action quality (accuracy) --
# each pair shares one color so the two metrics being tracked read apart. A
# translucent bubble (matching bar color) is drawn behind each pair, with a
# pill tab straddling its top edge naming which metric it is, so the two
# groups read as two visually separated cards, not one undifferentiated
# stack of bars.
VIDEO_QUALITY_COLOR = TEAL
ACTION_QUALITY_COLOR = PINK

bubble_left = chart_left - Inches(0.3)
bubble_width = chart_width + Inches(0.6)

# -- video-quality bubble (PSNR ladder) --------------------------------------
video_bubble_top = Inches(1.65)
video_bubble_bottom = Inches(4.50)
stage_band(s, bubble_left, video_bubble_top, bubble_width,
           video_bubble_bottom - video_bubble_top, VIDEO_QUALITY_COLOR, alpha_pct=18)
chip(s, Inches(1.0), video_bubble_top - Inches(0.25), Inches(4.1), Inches(0.5),
     "▶ VIDEO QUALITY", VIDEO_QUALITY_COLOR, 30)

rows = [
    ("before our training", 7.12, VIDEO_QUALITY_COLOR),
    ("after, with the real action", 18.56, VIDEO_QUALITY_COLOR),
]
top = Inches(2.15)
row_step = Inches(1.15)
bar_h = Inches(0.5)
for label, val, color in rows:
    add_text(s, chart_left, top, Inches(8), Inches(0.42), label, 30, DIM,
             bold=True)
    bar_y = top + Inches(0.45)
    rect(s, chart_left, bar_y, chart_width, bar_h, fill_color=PANEL, radius=True)
    bw = Emu(int(chart_width * (val / scale_max)))
    rect(s, chart_left, bar_y, bw, bar_h, fill_color=color, radius=True)
    add_text(s, chart_left + bw + Inches(0.15), bar_y - Inches(0.06),
             Inches(2.2), Inches(0.6), f"{val:.1f} dB", 30, INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    top += row_step

# delta callout, level with the "after" row's own label
add_text(s, Inches(9.15), top - row_step, Inches(2.2), Inches(0.5),
         "+11.4 dB", 36, AMBER, bold=True, align=PP_ALIGN.CENTER)

ceiling_x = chart_left + Emu(int(chart_width * (ceiling_db / scale_max)))
ceiling_top, ceiling_bottom = Inches(2.05), video_bubble_bottom - Inches(0.15)
conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ceiling_x, ceiling_top,
                              ceiling_x, ceiling_bottom)
conn.line.color.rgb = INK
conn.line.width = Pt(2.5)
conn.line.dash_style = MSO_LINE.DASH
add_text(s, ceiling_x - Inches(3.9), ceiling_top, Inches(4.0),
         Inches(0.45), "VAE ceiling · 22.7 dB", 30, INK, bold=True,
         align=PP_ALIGN.RIGHT)

# -- action-accuracy bubble (relative direction accuracy) --------------------
acc_bubble_top = Inches(4.85)
acc_bubble_bottom = Inches(7.25)
stage_band(s, bubble_left, acc_bubble_top, bubble_width,
           acc_bubble_bottom - acc_bubble_top, ACTION_QUALITY_COLOR, alpha_pct=18)
chip(s, Inches(1.0), acc_bubble_top - Inches(0.25), Inches(4.3), Inches(0.5),
     "▶ ACTION ACCURACY", ACTION_QUALITY_COLOR, 30)

acc_top = Inches(5.20)
acc_row_step = Inches(1.0)
acc_bar_h = Inches(0.4)
acc_scale = 100.0
acc_rows = [
    ("untrained", 51.5, ACTION_QUALITY_COLOR),
    ("fine-tuned", 99.6, ACTION_QUALITY_COLOR),
]
for label, val, color in acc_rows:
    add_text(s, chart_left, acc_top, Inches(3.0), Inches(0.42), label, 30,
             DIM, bold=True)
    bar_y = acc_top + Inches(0.45)
    rect(s, chart_left, bar_y, chart_width, acc_bar_h, fill_color=PANEL,
         radius=True)
    bw = Emu(int(chart_width * (val / acc_scale)))
    rect(s, chart_left, bar_y, bw, acc_bar_h, fill_color=color, radius=True)
    add_text(s, chart_left + bw + Inches(0.15), bar_y - Inches(0.13),
             Inches(1.6), Inches(0.6), f"{val:.1f}%", 30, INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    acc_top += acc_row_step

# ============================================ 7 · Interactive widget =======
s = new_slide(prs, RGBColor(0x10, 0x3C, 0x5C), RGBColor(0x2C, 0x1A, 0x52), 122)
add_text(s, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.75),
         "Live demo: click a command", 36, INK, bold=True)

scr = Inches(2.7)
sx, sy = Inches(0.7), Inches(1.5)
rect(s, sx - Inches(0.1), sy - Inches(0.1), scr + Inches(0.2), scr + Inches(0.2),
     fill_color=PANEL, line_color=CYAN, radius=True, line_w=2.0)
s.shapes.add_picture(os.path.join(WIDGET, "screen_context.png"), sx, sy, scr, scr)
add_text(s, sx, sy + scr + Inches(0.18), scr, Inches(0.55), "start frame",
         30, DIM, bold=True, align=PP_ALIGN.CENTER)

link = rect(s, sx - Inches(0.1), sy + scr + Inches(0.8), scr + Inches(0.2),
            Inches(0.75), fill_color=PANEL, line_color=TEAL, radius=True,
            line_w=2.25)
ltf = link.text_frame
ltf.word_wrap = False
ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
lp = ltf.paragraphs[0]
lp.alignment = PP_ALIGN.CENTER
lr = lp.add_run()
lr.text = "FULL DEMO ↗"
lr.font.size = Pt(30)
lr.font.bold = True
lr.font.name = FONT
lr.font.color.rgb = TEAL
link.click_action.hyperlink.address = "arm_control_panel.html"

# free-rollout tile dropped; the freed width goes into enlarging + spreading
# out the remaining five commands instead
tile = Inches(2.0)
xgap = Inches(2.35) - tile   # wide column spacing, freed by dropping the chain tile
ygap = Inches(0.075)         # rows are the tight dimension (3 stacked in slide height)
px = Inches(4.9)
py = Inches(1.25)
cols = [px, px + tile + xgap, px + 2 * (tile + xgap)]
rws = [py, py + tile + ygap, py + 2 * (tile + ygap)]
pad = [("up", cols[1], rws[0]), ("left", cols[0], rws[1]),
       ("still", cols[1], rws[1]), ("right", cols[2], rws[1]),
       ("down", cols[1], rws[2])]
for name, x, y in pad:
    s.shapes.add_movie(os.path.join(WIDGET, f"tile_{name}.mp4"), x, y, tile, tile,
                       poster_frame_image=os.path.join(WIDGET, f"tile_{name}.png"),
                       mime_type="video/mp4")

# ============================================= 8 · Inverse dynamics ========
# One slide per search step (was a single auto-playing GIF) -- the presenter
# now clicks forward through the search by hand instead of watching a timed
# loop. Frames come from make_search_tree_frames.py (same real data/argmin
# logic the old GIF used, just frozen into stills).
STEPS_DIR = "/home/mls10/presentation/search_steps"
N_STEPS = 7
for step_i in range(N_STEPS):
    s = new_slide(prs, RGBColor(0x14, 0x42, 0x4F), RGBColor(0x3B, 0x1E, 0x46), 118)
    eyebrow(s, Inches(0.8), Inches(0.45), "Inverse dynamics", VIOLET,
            width=Inches(6))
    add_text(s, Inches(10.3), Inches(0.5), Inches(2.5), Inches(0.4),
             f"click {step_i + 1} / {N_STEPS}", 30, DIM, bold=True,
             align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.8),
             "Choosing an action by imagining its outcome", 36, INK, bold=True)

    gif_w, gif_h = Inches(7.4), Inches(4.757)   # canvas 1120x720, exact 1.556 AR
    gx, gy = Inches(0.7), Inches(1.85)
    rect(s, gx - Inches(0.05), gy - Inches(0.05), gif_w + Inches(0.10),
         gif_h + Inches(0.10), fill_color=PANEL, line_color=VIOLET, radius=False,
         line_w=2.5)
    s.shapes.add_picture(os.path.join(STEPS_DIR, f"step{step_i}.png"), gx, gy,
                         gif_w, gif_h)

    rect(s, Inches(8.55), Inches(2.2), Inches(3.9), Inches(1.5), fill_color=PANEL,
         line_color=TEAL, radius=True, line_w=2.25)
    add_text(s, Inches(8.55), Inches(2.2), Inches(3.9), Inches(1.5), "15 / 20",
             48, TEAL, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.55), Inches(3.85), Inches(3.9), Inches(0.55),
             "imagination accuracy", 30, DIM, bold=True, align=PP_ALIGN.CENTER)

# ===================================================== 9 · Limitations ====
# Flat background (no diagonal gradient) -- unlike every other slide, this
# one carries three dense cards edge to edge, and the usual diagonal seam
# cut visibly across the bottom of them.
s = new_slide_solid(prs, RGBColor(0x2E, 0x17, 0x24))
eyebrow(s, Inches(0.8), Inches(0.5), "Free rollout", CORAL)
add_text(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.9),
         "Control holds. The scene drifts.", 40, INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1), Inches(3.2), CORAL, VIOLET)

card_w, card_h = Inches(3.75), Inches(2.6)
cgap = Inches(0.25)
cy = Inches(2.6)
cards = [
    ("98–100%", TEAL, "control holds"),
    ("79 → 72%", CORAL, "position drifts"),
    ("+97 → +57%", AMBER, "video quality decay"),
]
for i, (num, color, label) in enumerate(cards):
    x = Inches(0.8) + i * (card_w + cgap)
    rect(s, x, cy, card_w, card_h, fill_color=PANEL, line_color=color,
         radius=True, line_w=2.5)
    add_text(s, x, cy + Inches(0.35), card_w, Inches(0.9), num, 44, color,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), cy + Inches(1.45), card_w - Inches(0.4),
             Inches(1.0), label, 30, DIM, bold=True, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.9),
         "A known cause — and a known fix",
         32, INK, bold=True)

# ================================================== 10 · What comes next ====
s = new_slide(prs, RGBColor(0x1F, 0x1C, 0x56), RGBColor(0x0F, 0x46, 0x40), 125)
eyebrow(s, Inches(0.8), Inches(0.5), "What's next", AMBER)

add_text(s, Inches(0.6), Inches(1.35), Inches(5.6), Inches(1.7), "5 days", 100,
         AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.2), Inches(5.6), Inches(0.55),
         "us, on one A100 GPU", 30, INK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.2), Inches(1.9), Inches(1.0), Inches(0.8), "vs", 40, DIM,
         bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.2), Inches(1.35), Inches(5.6), Inches(1.7), "months", 100,
         SLATE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.2), Inches(3.2), Inches(5.6), Inches(0.55),
         "typical for this problem", 30, DIM, align=PP_ALIGN.CENTER)

accent_bar(s, Inches(0.8), Inches(4.25), Inches(11.7), AMBER, TEAL)

steps = [
    ("scene anchoring, so long rollouts stop drifting", TEAL),
    ("true self-forcing training", CYAN),
    ("256 × 256: lifts the 22.7 dB fidelity ceiling", VIOLET),
]
top = Inches(4.7)
for txt, c in steps:
    add_text(s, Inches(0.9), top, Inches(0.6), Inches(0.55), "→", 30, c, bold=True)
    add_text(s, Inches(1.6), top, Inches(11.2), Inches(0.55), txt, 30, INK, bold=True)
    top += Inches(0.85)

# ============================================================ 11 · Q&A ====
s = new_slide(prs, RGBColor(0x1A, 0x1F, 0x5A), RGBColor(0x0D, 0x44, 0x4A), 128)
add_text(s, Inches(0), Inches(0.5), Inches(13.333), Inches(1.3), "Q&A", 84,
         INK, bold=True, align=PP_ALIGN.CENTER)
accent_bar(s, Inches(5.77), Inches(1.95), Inches(1.8), PINK, TEAL, Pt(5))

# real headshots (center-cropped square, prep_qa_assets.py) + a scannable QR
# per person so the audience can open LinkedIn with a phone camera, not just
# a click -- the whole card stays hyperlinked too, for anyone viewing on
# screen.
people = [
    ("Mihajlo Stevanović", "mihastevanovic04@gmail.com",
     "https://www.linkedin.com/in/mihajlo-stevanovi%C4%87/",
     "/home/mls10/presentation/qa_mihajlo_sq.jpg",
     "/home/mls10/presentation/qa_qr_mihajlo.png", CYAN),
    ("David Marković", "dawidzards@gmail.com",
     "https://www.linkedin.com/in/david-markovic-3a107a309/",
     "/home/mls10/presentation/qa_david_sq.jpg",
     "/home/mls10/presentation/qa_qr_david.png", PINK),
]
photo_size = Inches(1.9)
frame_pad = Inches(0.05)
for i, (name, mail, url, photo, qr, col) in enumerate(people):
    x = Inches(0.65) + i * Inches(6.3)
    card = rect(s, x, Inches(2.3), Inches(5.7), Inches(4.5), fill_color=PANEL,
                line_color=col, radius=True, line_w=2.5)
    card.click_action.hyperlink.address = url

    add_text(s, x + Inches(0.2), Inches(2.55), Inches(5.3), Inches(0.55), name,
             32, INK, bold=True, align=PP_ALIGN.CENTER)

    px, py = x + Inches(0.5), Inches(3.25)
    rect(s, px - frame_pad, py - frame_pad, photo_size + 2 * frame_pad,
         photo_size + 2 * frame_pad, fill_color=PANEL, line_color=col,
         line_w=2.5)
    s.shapes.add_picture(photo, px, py, photo_size, photo_size)

    qx, qy = x + Inches(3.3), Inches(3.25)
    rect(s, qx - frame_pad, qy - frame_pad, photo_size + 2 * frame_pad,
         photo_size + 2 * frame_pad, fill_color=PANEL, line_color=col,
         line_w=2.5)
    s.shapes.add_picture(qr, qx, qy, photo_size, photo_size)

    add_text(s, x + Inches(0.2), Inches(5.35), Inches(5.3), Inches(0.5),
             "scan for LinkedIn ↗", 30, col, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(5.95), Inches(5.3), Inches(0.5), mail,
             30, DIM, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print("saved:", OUT_PATH, f"{os.path.getsize(OUT_PATH)/1e6:.2f} MB")
